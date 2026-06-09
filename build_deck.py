"""Builds the Liquify and Container World presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1F, 0x32, 0x4E)
STEEL = RGBColor(0x3E, 0x5C, 0x76)
LIGHT = RGBColor(0xEE, 0xF2, 0xF6)
GREY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x28, 0x33)
FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

DECK_TITLE = "Liquify and Container World: Order to Cash"


def add_slide():
    return prs.slides.add_slide(blank)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, size, color=DARK,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_header(slide, title, number):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    set_fill(bar, NAVY)
    add_text(slide, Inches(0.6), Inches(0.35), Inches(11.0), Inches(0.8),
             title, 30, NAVY, bold=True)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.08),
                                  Inches(1.6), Inches(0.045))
    set_fill(rule, STEEL)
    add_text(slide, Inches(0.6), Inches(7.02), Inches(9.5), Inches(0.35),
             DECK_TITLE, 10, GREY)
    add_text(slide, Inches(12.3), Inches(7.02), Inches(0.45), Inches(0.35),
             str(number), 10, GREY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, top=Inches(1.55), left=Inches(0.65),
                width=Inches(12.0), size=18, space_after=14):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
        pPr.set("marL", "274320")
        pPr.set("indent", "-274320")
        r1 = p.add_run()
        r1.text = "▪  "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = STEEL
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = DARK
    return box


def flow_box(slide, left, top, width, height, number, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.08
    set_fill(shape, LIGHT)
    shape.line.color.rgb = STEEL
    shape.line.width = Pt(1.0)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rn = p.add_run()
    rn.text = f"{number}. "
    rn.font.name = FONT
    rn.font.size = Pt(13)
    rn.font.bold = True
    rn.font.color.rgb = STEEL
    rt = p.add_run()
    rt.text = text
    rt.font.name = FONT
    rt.font.size = Pt(13)
    rt.font.color.rgb = DARK
    return shape


def right_arrow(slide, left, top, width, height):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    set_fill(arr, STEEL)
    arr.shadow.inherit = False
    return arr


def elbow_arrow(slide, begin_shape, end_shape):
    """Elbow connector with an arrowhead from bottom of begin to top of end."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, 0, 0, Emu(1), Emu(1))
    conn.begin_connect(begin_shape, 3)  # bottom
    conn.end_connect(end_shape, 0)      # top
    conn.line.color.rgb = STEEL
    conn.line.width = Pt(2.25)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def flow_chart(slide, steps, row_split):
    """Two rows of boxes, left to right, joined by an elbow arrow at the wrap."""
    row1 = steps[:row_split]
    row2 = steps[row_split:]
    arrow_w = Inches(0.42)
    arrow_h = Inches(0.32)
    margin = Inches(0.65)
    avail = SLIDE_W - margin * 2

    def layout_row(items, top, box_h):
        n = len(items)
        box_w = int((avail - arrow_w * (n - 1)) / n)
        shapes = []
        x = margin
        for i, text in enumerate(items):
            s = flow_box(slide, x, top, box_w, box_h, text[0], text[1])
            shapes.append(s)
            x += box_w
            if i < n - 1:
                right_arrow(slide, x, int(top + box_h / 2 - arrow_h / 2),
                            arrow_w, arrow_h)
                x += arrow_w
        return shapes

    box_h = Inches(1.55)
    r1_shapes = layout_row(row1, Inches(2.0), box_h)
    r2_shapes = layout_row(row2, Inches(4.55), box_h)
    elbow_arrow(slide, r1_shapes[-1], r2_shapes[0])


# Slide 1
s = add_slide()
add_text(s, Inches(0.6), Inches(0.32), Inches(12.0), Inches(0.35),
         DECK_TITLE.upper(), 13, STEEL, bold=True)
add_text(s, Inches(0.6), Inches(0.72), Inches(11.0), Inches(0.8),
         "What Liquify Is", 30, NAVY, bold=True)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.45),
                          Inches(1.6), Inches(0.045))
set_fill(rule, STEEL)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
set_fill(bar, NAVY)
add_text(s, Inches(0.6), Inches(7.02), Inches(9.5), Inches(0.35), DECK_TITLE, 10, GREY)
add_text(s, Inches(12.3), Inches(7.02), Inches(0.45), Inches(0.35), "1", 10, GREY,
         align=PP_ALIGN.RIGHT)
add_bullets(s, [
    "Liquify is a private B2B marketplace for the BC beverage alcohol trade, founded 2014 and based in Victoria.",
    "It connects agents and suppliers with over 3,000 BC licensees across retail and on premise.",
    "Buyers browse a catalogue of 27,000+ LDB and Direct Delivery products and order on platform instead of by phone or email.",
    "For us, it converts a rep's pitch into a confirmed, legally binding order on the spot.",
], top=Inches(1.95), size=19, space_after=18)

# Slide 2
s = add_slide()
add_header(s, "How It Integrates With Container World", 2)
add_bullets(s, [
    "Our Direct Delivery inventory is warehoused at Container World in Richmond.",
    "When a licensee places an order on Liquify and we confirm it on the platform, Liquify pushes the order directly into Container World's system.",
    "Container World picks the order and delivers province wide through its trucking arm, Commercial Logistics Inc.",
    "Liquify tracks order status from dispatch through fulfillment, so confirmation, fulfillment, and delivery status all live in one place.",
    "No rekeying. No order desk emails.",
], size=19, space_after=18)

# Slide 3
s = add_slide()
add_header(s, "Order Flow", 3)
flow_chart(s, [
    (1, "Licensee browses catalogue and builds Order Sheet"),
    (2, "Licensee submits order on Liquify"),
    (3, "We review and confirm on platform. Order becomes a binding agreement."),
    (4, "Liquify pushes confirmed order to Container World, Richmond"),
    (5, "Container World picks and stages"),
    (6, "Commercial Logistics Inc. delivers to licensee"),
    (7, "Delivery status updates in Liquify"),
], row_split=4)

# Slide 4
s = add_slide()
add_header(s, "Payment Flow", 4)
flow_chart(s, [
    (1, "Order confirmed and fulfilled"),
    (2, "Liquify auto generates the Doc 60 invoice with the licensee's details and delivers it to us"),
    (3, "Licensee's credit card on file is charged under their stored card authorization"),
    (4, "Funds settle through a PCI DSS compliant third party processor into our own merchant account"),
    (5, "We remit LDB markup, excise, GST, and PST"),
    (6, "Doc 60 and transaction records stored in Liquify for reconciliation"),
], row_split=3)

# Slide 5
s = add_slide()
add_header(s, "What Finance Needs to Know", 5)
add_bullets(s, [
    "Liquify is a venue, not a payment company. It never takes ownership of product or custody of funds.",
    "We hold the merchant account, so card processing fees and any chargeback exposure are ours.",
    "Liquify bills us a transaction fee when each Doc 60 is issued, charged to our card on file, plus an annual platform fee.",
    "Tax remittance stays with us.",
    "Net effect: faster confirmed orders, automated invoicing, and payment collected by card authorization instead of chasing receivables on terms.",
], size=19, space_after=18)

OUT = "Liquify and Container World Order to Cash.pptx"
prs.save(OUT)
print("saved", OUT)
